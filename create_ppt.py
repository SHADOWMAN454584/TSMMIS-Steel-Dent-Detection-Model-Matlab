"""
Generate a comprehensive PowerPoint presentation for the TSMMIS MATLAB project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──────────────────────────────────────────────────────────────
BASE = r"D:\MVRP_Project\TSMMIS_MATLAB"
TSNE_IMG = os.path.join(BASE, "evaluation", "results", "tsne_visualization.png")
DATA_DIR = os.path.join(BASE, "data", "NEU-CLS")
OUTPUT   = os.path.join(BASE, "TSMMIS_Presentation.pptx")

# Sample images – one per defect class
SAMPLE_IMGS = {
    "Crakling":       os.path.join(DATA_DIR, "cra", "crakling_1.jpg"),
    "Inclusion":      os.path.join(DATA_DIR, "in", "inclusion_1.jpg"),
    "Patches":        os.path.join(DATA_DIR, "pa", "patches_1.jpg"),
    "Pitted Surface": os.path.join(DATA_DIR, "ps", "pitted_surface_1.jpg"),
    "Rolled-in Scale":os.path.join(DATA_DIR, "rs", "rolled-in_scale_1.jpg"),
    "Scratches":      os.path.join(DATA_DIR, "sc", "scratches_1.jpg"),
}

# ── Color palette ──────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)   # Deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x96, 0xD6)   # Bright blue
ACCENT_TEAL  = RGBColor(0x00, 0xB4, 0xA0)   # Teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
GOLD         = RGBColor(0xFF, 0xD7, 0x00)
CARD_BG      = RGBColor(0x24, 0x24, 0x3E)   # Slightly lighter navy
TABLE_HDR    = RGBColor(0x00, 0x70, 0xA0)
TABLE_ROW1   = RGBColor(0x1E, 0x1E, 0x38)
TABLE_ROW2   = RGBColor(0x28, 0x28, 0x45)
RED_ACCENT   = RGBColor(0xE8, 0x4D, 0x4D)
GREEN_ACCENT = RGBColor(0x4E, 0xCB, 0x71)
ORANGE       = RGBColor(0xFF, 0xA5, 0x00)
PURPLE       = RGBColor(0xA8, 0x5C, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


# ── Helper functions ───────────────────────────────────────────────────
def add_bg(slide, color=DARK_BG):
    """Fill the slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with single-run text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, items, font_size=16,
                              color=WHITE, bullet_color=ACCENT_TEAL, font_name="Calibri"):
    """Add multi-line bulleted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        # Bullet character
        run_b = p.add_run()
        run_b.text = "\u25B8 "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = font_name
        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = font_name
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """Add a thin accent line (actually a very thin rectangle)."""
    return add_shape(slide, left, top, width, height, fill_color=color)


def add_section_title(slide, title_text, subtitle_text=None, number=None):
    """Standard section header at top of slide."""
    y = Inches(0.4)
    if number is not None:
        # Number badge
        badge = add_rounded_rect(slide, Inches(0.6), y, Inches(0.55), Inches(0.55), fill_color=ACCENT_BLUE)
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge.text_frame.paragraphs[0].add_run()
        run.text = str(number)
        run.font.size = Pt(22)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Calibri"
        badge.text_frame.paragraphs[0].space_before = Pt(4)
        title_left = Inches(1.35)
    else:
        title_left = Inches(0.6)

    add_text_box(slide, title_left, y - Inches(0.05), Inches(10), Inches(0.55),
                 title_text, font_size=30, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), y + Inches(0.55), Inches(2.5), ACCENT_BLUE)

    if subtitle_text:
        add_text_box(slide, Inches(0.6), y + Inches(0.65), Inches(11), Inches(0.4),
                     subtitle_text, font_size=15, color=LIGHT_GRAY)


def add_table(slide, left, top, rows_data, col_widths, font_size=13):
    """Add a styled table.  rows_data[0] = header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_width = sum(col_widths)
    row_height = Inches(0.42)
    table = slide.shapes.add_table(n_rows, n_cols, left, top,
                                   table_width, row_height * n_rows).table
    for c, w in enumerate(col_widths):
        table.columns[c].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HDR
            else:
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW1 if r % 2 == 1 else TABLE_ROW2
    return table


# =====================================================================
# SLIDE 1 – TITLE SLIDE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Top accent bar
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)

# Central content
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.5),
             "TSMMIS", font_size=56, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(0.9),
             "Teacher-Student Multi-Modal Instance Similarity",
             font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(3.3), Inches(3.3), ACCENT_TEAL, Pt(3))

add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.6),
             "Self-Supervised Few-Shot Steel Surface Defect Recognition",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.5),
             "Implemented in MATLAB  |  NEU-CLS Dataset  |  ResNet-18 Backbone",
             font_size=16, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

# Bottom line
add_shape(slide, 0, SH - Inches(0.06), SW, Inches(0.06), fill_color=ACCENT_BLUE)


# =====================================================================
# SLIDE 2 – TABLE OF CONTENTS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Presentation Outline")

toc_items = [
    ("01", "Introduction & Problem Statement"),
    ("02", "Dataset – NEU-CLS Steel Surface Defects"),
    ("03", "System Architecture & Pipeline"),
    ("04", "Data Preprocessing & Augmentation"),
    ("05", "Self-Supervised Pretraining"),
    ("06", "Loss Functions (Contrastive + MMIS)"),
    ("07", "Few-Shot Classification"),
    ("08", "Prototype Self-Training Adaptation"),
    ("09", "Hyperparameters"),
    ("10", "Experimental Results"),
    ("11", "t-SNE Feature Visualization"),
    ("12", "Project Structure"),
    ("13", "Conclusion"),
]

cols = 2
per_col = 7
for idx, (num, title) in enumerate(toc_items):
    col = idx // per_col
    row = idx % per_col
    x = Inches(1.0) + col * Inches(6.0)
    y = Inches(1.6) + row * Inches(0.72)
    # Number box
    nb = add_rounded_rect(slide, x, y, Inches(0.5), Inches(0.45), fill_color=ACCENT_BLUE)
    nb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = nb.text_frame.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Calibri"
    # Title
    add_text_box(slide, x + Inches(0.65), y + Inches(0.02), Inches(5), Inches(0.42),
                 title, font_size=17, color=WHITE)


# =====================================================================
# SLIDE 3 – INTRODUCTION & PROBLEM STATEMENT
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Introduction & Problem Statement", number=1)

# Problem card
card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=CARD_BG)
add_text_box(slide, Inches(1.0), Inches(1.65), Inches(5), Inches(0.4),
             "THE PROBLEM", font_size=18, color=RED_ACCENT, bold=True)
add_bullet_slide_content(slide, Inches(1.0), Inches(2.15), Inches(5), Inches(4.3), [
    "Steel manufacturing requires rigorous surface quality inspection",
    "Manual inspection is slow, subjective, and error-prone",
    "Deep learning needs thousands of labeled images per defect class",
    "Labeling steel defect images is expensive and requires domain experts",
    "In real-world scenarios, only 1-4 labeled samples per class are available",
], font_size=15, bullet_color=RED_ACCENT)

# Solution card
card2 = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=CARD_BG)
add_text_box(slide, Inches(7.3), Inches(1.65), Inches(5), Inches(0.4),
             "OUR SOLUTION: TSMMIS", font_size=18, color=GREEN_ACCENT, bold=True)
add_bullet_slide_content(slide, Inches(7.3), Inches(2.15), Inches(5), Inches(4.3), [
    "Self-supervised contrastive learning from unlabeled images",
    "Teacher-Student framework for robust feature learning",
    "Multi-Modal Instance Similarity (MMIS) regularization",
    "Few-shot classification with only 1-4 labeled examples per class",
    "Prototype-based classifier with self-training adaptation",
    "Achieves up to 95.75% accuracy with just 4 labeled samples/class",
], font_size=15, bullet_color=GREEN_ACCENT)


# =====================================================================
# SLIDE 4 – DATASET  (with sample images)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Dataset – NEU-CLS Steel Surface Defects", number=2)

add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             "Northeastern University Classification Dataset  |  6 defect classes  |  1,440 grayscale images  |  200x200 pixels",
             font_size=15, color=LIGHT_GRAY)

# Place 6 sample images in a row
img_size = Inches(1.65)
gap = Inches(0.25)
total_w = 6 * img_size + 5 * gap
start_x = (SW - total_w) / 2
img_y = Inches(1.85)

class_names = ["Crakling", "Inclusion", "Patches", "Pitted Surface", "Rolled-in Scale", "Scratches"]
class_colors = [RED_ACCENT, ORANGE, GOLD, ACCENT_TEAL, ACCENT_BLUE, PURPLE]
img_keys = list(SAMPLE_IMGS.keys())

for i, cls_name in enumerate(img_keys):
    x = start_x + i * (img_size + gap)
    img_path = SAMPLE_IMGS[cls_name]
    # Card behind image
    add_rounded_rect(slide, x - Inches(0.08), img_y - Inches(0.08),
                     img_size + Inches(0.16), img_size + Inches(0.55),
                     fill_color=CARD_BG)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, img_y, img_size, img_size)
    # Label
    add_text_box(slide, x - Inches(0.08), img_y + img_size + Inches(0.05),
                 img_size + Inches(0.16), Inches(0.35),
                 cls_name, font_size=12, color=class_colors[i], bold=True,
                 alignment=PP_ALIGN.CENTER)

# Dataset statistics table
table_data = [
    ["Class", "Abbreviation", "Train (60%)", "Test (40%)", "Total"],
    ["Crakling", "cra", "144", "96", "240"],
    ["Inclusion", "in", "144", "96", "240"],
    ["Patches", "pa", "144", "96", "240"],
    ["Pitted Surface", "ps", "144", "96", "240"],
    ["Rolled-in Scale", "rs", "144", "96", "240"],
    ["Scratches", "sc", "144", "96", "240"],
    ["TOTAL", "---", "864", "576", "1,440"],
]
col_widths = [Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
add_table(slide, Inches(1.8), Inches(4.3), table_data, col_widths, font_size=12)


# =====================================================================
# SLIDE 5 – SYSTEM ARCHITECTURE & PIPELINE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "System Architecture & Pipeline", number=3)

steps = [
    ("1", "Data\nPreparation", "Load NEU-CLS\n60/40 split\nAugmentation", ACCENT_BLUE),
    ("2", "Model\nInit", "ResNet-18\nImageNet pretrained\n512-dim features", ACCENT_TEAL),
    ("3", "Self-Supervised\nPretraining", "PCA projection\nK-means clustering\n300-epoch refinement", PURPLE),
    ("4", "Few-Shot\nEvaluation", "N=1,2,3,4 shots\n100 random trials\nPrototype classifier", ORANGE),
    ("5", "Prototype\nAdaptation", "Self-training\nConfidence margin\nUnlabeled fusion", RED_ACCENT),
    ("6", "Results &\nVisualization", "Accuracy table\nt-SNE embedding\nScatter plot", GREEN_ACCENT),
]

box_w = Inches(1.7)
box_h = Inches(2.8)
gap = Inches(0.22)
total = len(steps) * box_w + (len(steps) - 1) * gap
sx = (SW - total) / 2
sy = Inches(1.7)

for i, (num, title, desc, color) in enumerate(steps):
    x = sx + i * (box_w + gap)
    # Card
    card = add_rounded_rect(slide, x, sy, box_w, box_h, fill_color=CARD_BG)
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + box_w/2 - Inches(0.25), sy + Inches(0.2),
                                   Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = circ.text_frame.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(18)
    r.font.color.rgb = WHITE
    r.font.bold = True
    # Title
    add_text_box(slide, x + Inches(0.1), sy + Inches(0.85), box_w - Inches(0.2), Inches(0.7),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, x + Inches(0.1), sy + Inches(1.55) + j * Inches(0.3),
                     box_w - Inches(0.2), Inches(0.3),
                     line, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(steps) - 1:
        add_text_box(slide, x + box_w, sy + box_h/2 - Inches(0.15), gap, Inches(0.3),
                     "\u25B6", font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.8),
             "Pipeline is orchestrated by main_TSMMIS.m (683 lines) which calls modular components:\n"
             "DataLoader.m | ResNet18Encoder.m | ContrastiveLoss.m | MMISLoss.m | TSMMISTrainer.m | TSMMISEvaluator.m",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 6 – DATA PREPROCESSING & AUGMENTATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Data Preprocessing & Augmentation", number=4)

# Left panel - preprocessing
add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=CARD_BG)
add_text_box(slide, Inches(1.0), Inches(1.65), Inches(5), Inches(0.4),
             "PREPROCESSING", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(1.0), Inches(2.2), Inches(5), Inches(2.5), [
    "Images resized from 200x200 to 224x224 pixels",
    "Per-class folder scanning (JPG/PNG/BMP supported)",
    "Random 60/40 train-test split per class",
    "From 144 train images: 4 reserved few-shot + 140 unlabeled",
    "Full training pool used for N-shot sampling",
], font_size=14)

# Right panel - augmentation
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), fill_color=CARD_BG)
add_text_box(slide, Inches(7.3), Inches(1.65), Inches(5), Inches(0.4),
             "DATA AUGMENTATION (Multi-Crop)", font_size=18, color=ACCENT_TEAL, bold=True)
add_bullet_slide_content(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(2.5), [
    "K = 5 multi-crop views per image",
    "Crop ratio range: [0.1, 1.0]",
    "Random horizontal & vertical flipping",
    "Color jitter: brightness +/-20%, contrast +/-20%",
    "Random grayscale conversion (20% probability)",
    "Gaussian blur (30% probability)",
], font_size=14, bullet_color=ACCENT_TEAL)

# Bottom flow diagram
add_text_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
             "Raw Image (200x200) --> Resize (224x224) --> Multi-Crop (K=5) --> Augment (flip, jitter, blur) --> Feature Extraction",
             font_size=14, color=GOLD, alignment=PP_ALIGN.CENTER, bold=True)

add_text_box(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
             "Implemented in: utils/DataLoader.m (284 lines)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 7 – SELF-SUPERVISED PRETRAINING
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Self-Supervised Pretraining", number=5)

pretrain_steps = [
    ("Step 1", "Feature Extraction", "Extract 512-dim features from\nall unlabeled images using\nResNet-18 encoder (batch=32)", ACCENT_BLUE),
    ("Step 2", "L2 Normalization", "Row-wise L2 normalization\nof all feature vectors", ACCENT_TEAL),
    ("Step 3", "PCA Projection", "Project to 256 dimensions\nwith centering and whitening\nthen L2-normalize", PURPLE),
    ("Step 4", "K-Means Init", "Initialize 6 cluster centers\nwith 5 replicates\nfor robust seeding", ORANGE),
    ("Step 5", "Prototype Refinement", "300 epochs: cosine similarity\nassignment + mean re-computation\nTrack loss = mean(1 - maxSim)", RED_ACCENT),
]

box_w = Inches(2.2)
box_h = Inches(2.6)
gap_x = Inches(0.2)
total_w = len(pretrain_steps) * box_w + (len(pretrain_steps) - 1) * gap_x
sx = (SW - total_w) / 2
sy = Inches(1.6)

for i, (step, title, desc, color) in enumerate(pretrain_steps):
    x = sx + i * (box_w + gap_x)
    add_rounded_rect(slide, x, sy, box_w, box_h, fill_color=CARD_BG)
    add_text_box(slide, x, sy + Inches(0.15), box_w, Inches(0.3),
                 step, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, sy + Inches(0.45), box_w, Inches(0.4),
                 title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(0.3), sy + Inches(0.9), box_w - Inches(0.6), color, Pt(2))
    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, x + Inches(0.15), sy + Inches(1.1) + j * Inches(0.32),
                     box_w - Inches(0.3), Inches(0.32),
                     line, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(pretrain_steps) - 1:
        add_text_box(slide, x + box_w, sy + box_h / 2 - Inches(0.15), gap_x, Inches(0.3),
                     "\u25B6", font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Output
add_rounded_rect(slide, Inches(2), Inches(4.55), Inches(9.3), Inches(1.1), fill_color=CARD_BG)
add_text_box(slide, Inches(2.3), Inches(4.65), Inches(8.7), Inches(0.35),
             "OUTPUT: Learned Transform", font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.3), Inches(5.05), Inches(8.7), Inches(0.45),
             "{ mu (mean), coeff (PCA matrix), prototypeCenters (6 class prototypes), explainedVariance }",
             font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.35),
             "Pretraining uses no labels -- learns structure from unlabeled data alone",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 8 – LOSS FUNCTIONS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Loss Functions", number=6)

# Contrastive Loss card
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), fill_color=CARD_BG)
add_text_box(slide, Inches(0.9), Inches(1.65), Inches(5.2), Inches(0.4),
             "Contrastive Loss  L_t&s", font_size=20, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.3),
             "Teacher-Student Cross-View Contrastive Loss", font_size=13, color=LIGHT_GRAY)
add_accent_line(slide, Inches(0.9), Inches(2.5), Inches(5.2), ACCENT_BLUE, Pt(2))

equations_l = [
    "sim_ps_yk = (P_s . Y_k^T) / tau",
    "sim_pk_ys = (P_k . Y_s^T) / tau",
    "",
    "L_t&s = [ CE(sim_ps_yk, labels)",
    "        + CE(sim_pk_ys, labels) ] / (2K)",
]
for j, eq in enumerate(equations_l):
    add_text_box(slide, Inches(1.2), Inches(2.7) + j * Inches(0.35), Inches(4.8), Inches(0.35),
                 eq, font_size=14, color=GOLD, font_name="Consolas")

add_bullet_slide_content(slide, Inches(0.9), Inches(4.7), Inches(5.2), Inches(2.0), [
    "P_s = predictor output from main view",
    "Y_k = encoder outputs from K multi-crop views",
    "Y_s = main view features (detached for gradient)",
    "tau = 0.1 (temperature parameter)",
    "Cross-entropy with softmax over similarity rows",
], font_size=12, bullet_color=ACCENT_BLUE)

# MMIS Loss card
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), fill_color=CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.65), Inches(5.2), Inches(0.4),
             "MMIS Regularization  L_MMIS", font_size=20, color=ACCENT_TEAL, bold=True)
add_text_box(slide, Inches(7.2), Inches(2.1), Inches(5.2), Inches(0.3),
             "Multi-Modal Instance Similarity Loss", font_size=13, color=LIGHT_GRAY)
add_accent_line(slide, Inches(7.2), Inches(2.5), Inches(5.2), ACCENT_TEAL, Pt(2))

equations_r = [
    "M_k = norm(Y_k) . norm(Y_{k+1})^T",
    "pseudo_labels = argmax(M_k, dim=2)",
    "ce_k = CrossEntropy(M_k, pseudo_labels)",
    "L_MMIS = mean(-ce_k) over K-1 pairs",
    "L_total = lambda * epoch * L_MMIS",
]
for j, eq in enumerate(equations_r):
    add_text_box(slide, Inches(7.5), Inches(2.7) + j * Inches(0.35), Inches(4.8), Inches(0.35),
                 eq, font_size=14, color=GOLD, font_name="Consolas")

add_bullet_slide_content(slide, Inches(7.2), Inches(4.7), Inches(5.2), Inches(2.0), [
    "Adjacent multi-crop pair similarity matrices",
    "Pseudo-labels via per-row argmax",
    "Negative cross-entropy encourages diversity",
    "lambda = 0.001 (MMIS weight)",
    "Epoch-scaled: grows importance during training",
], font_size=12, bullet_color=ACCENT_TEAL)

# File references
add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.3),
             "Implemented in: training/ContrastiveLoss.m (99 lines)  |  training/MMISLoss.m (141 lines)",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 9 – FEW-SHOT CLASSIFICATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Few-Shot Classification", number=7)

# Left - Prototype Classifier
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), fill_color=CARD_BG)
add_text_box(slide, Inches(0.9), Inches(1.65), Inches(5.2), Inches(0.4),
             "Prototype Classifier", font_size=20, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(0.9), Inches(2.15), Inches(3.5), ACCENT_BLUE, Pt(2))

proto_eqs = [
    "prototype_c = L2_norm( mean(",
    "  support_features where label = c ) )",
    "",
    "proto_sim = features . prototypes^T",
    "support_sim = max( features . ",
    "  support_features_c^T )  per class c",
    "",
    "combined = 0.75 * proto_sim",
    "         + 0.25 * support_sim",
    "prediction = argmax( combined )",
]
for j, eq in enumerate(proto_eqs):
    add_text_box(slide, Inches(1.2), Inches(2.4) + j * Inches(0.32), Inches(5), Inches(0.32),
                 eq, font_size=13, color=GOLD, font_name="Consolas")

# Right - Evaluation Protocol
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), fill_color=CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.65), Inches(5.2), Inches(0.4),
             "Evaluation Protocol", font_size=20, color=ACCENT_TEAL, bold=True)
add_accent_line(slide, Inches(7.2), Inches(2.15), Inches(3.5), ACCENT_TEAL, Pt(2))

add_bullet_slide_content(slide, Inches(7.2), Inches(2.4), Inches(5.2), Inches(4.0), [
    "Label counts tested: N = 1, 2, 3, 4 per class",
    "100 random trials per label-count setting",
    "Each trial: random N-shot sampling from train pool",
    "Compute per-class prototype from support features",
    "Adapt prototypes with unlabeled data (self-training)",
    "Predict on test set using combined similarity scoring",
    "Record accuracy; report mean +/- std across 100 trials",
    "Blend ratio: 75% prototype + 25% nearest-instance",
], font_size=13, bullet_color=ACCENT_TEAL)


# =====================================================================
# SLIDE 10 – PROTOTYPE SELF-TRAINING ADAPTATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Prototype Self-Training Adaptation", number=8)

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
             "Leveraging unlabeled data to refine class prototypes through confident pseudo-label predictions",
             font_size=15, color=LIGHT_GRAY)

# Algorithm box
add_rounded_rect(slide, Inches(1.5), Inches(1.9), Inches(10.3), Inches(4.2), fill_color=CARD_BG)
add_text_box(slide, Inches(1.9), Inches(2.05), Inches(6), Inches(0.4),
             "ALGORITHM: Self-Training Adaptation", font_size=18, color=GOLD, bold=True)
add_accent_line(slide, Inches(1.9), Inches(2.5), Inches(9.5), GOLD, Pt(2))

algo_lines = [
    ("FOR", "iter = 1 to 2:"),
    ("  1.", "Compute similarity = unlabeled_features . prototypes^T"),
    ("  2.", "Sort similarity scores: get top-1 and top-2 per sample"),
    ("  3.", "Calculate confidence margin = top1_score - top2_score"),
    ("  4.", "IF margin >= 0.08 (high confidence):"),
    ("    ", "    pseudo_label = argmax(similarity)"),
    ("  5.", "Update prototype for each class c:"),
    ("    ", "    prototype_c = L2_norm( 0.70 * prototype_c + 0.30 * mean(confident_unlabeled_c) )"),
    ("END", ""),
]

for j, (prefix, text) in enumerate(algo_lines):
    y = Inches(2.7) + j * Inches(0.35)
    add_text_box(slide, Inches(2.2), y, Inches(1.2), Inches(0.35),
                 prefix, font_size=13, color=ACCENT_BLUE, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(3.0), y, Inches(8.3), Inches(0.35),
                 text, font_size=13, color=WHITE, font_name="Consolas")

# Key parameters
add_rounded_rect(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.7), fill_color=CARD_BG)
add_text_box(slide, Inches(1.9), Inches(6.4), Inches(9.5), Inches(0.4),
             "Key Parameters:  Iterations = 2  |  Confidence Margin Threshold = 0.08  |  Adaptation Blend = 0.70/0.30",
             font_size=14, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 11 – HYPERPARAMETERS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Hyperparameters", number=9)

# Table
hyp_data = [
    ["Parameter", "Value", "Description"],
    ["Image Size", "224 x 224", "Input resolution after resize"],
    ["K (Multi-Crop Views)", "5", "Number of augmented views per image"],
    ["Crop Ratio Range", "[0.1, 1.0]", "Min/max crop scale factor"],
    ["Train/Test Split", "60% / 40%", "Per-class data division"],
    ["Pretrain Epochs", "300", "Prototype refinement iterations"],
    ["Batch Size", "64", "Training batch size"],
    ["Learning Rate", "0.003", "SGD learning rate"],
    ["SGD Momentum", "0.9", "Optimizer momentum"],
    ["Weight Decay", "1e-4", "L2 regularization"],
    ["Lambda (MMIS)", "0.001", "MMIS loss weight"],
    ["Tau (Temperature)", "0.1", "Contrastive loss temperature"],
    ["PCA Dimension", "256", "Projected feature dimension"],
    ["Few-Shot Runs", "100", "Random trials per setting"],
    ["Label Counts", "[1, 2, 3, 4]", "Shots per class tested"],
    ["Self-Train Iterations", "2", "Adaptation loop count"],
    ["Confidence Margin", "0.08", "Pseudo-label threshold"],
    ["Proto/Support Blend", "0.75 / 0.25", "Combined similarity weights"],
    ["Adaptation Blend", "0.70 / 0.30", "Old/new prototype weights"],
]
col_widths_h = [Inches(3.0), Inches(2.0), Inches(4.5)]
add_table(slide, Inches(1.9), Inches(1.3), hyp_data, col_widths_h, font_size=11)


# =====================================================================
# SLIDE 12 – EXPERIMENTAL RESULTS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Experimental Results", number=10)

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
             "Few-Shot Classification Accuracy on NEU-CLS Test Set (576 images, 100 random trials per setting)",
             font_size=15, color=LIGHT_GRAY)

# Results table
res_data = [
    ["Setting", "Labels/Class", "Mean Accuracy", "Std Dev", "Interpretation"],
    ["1-shot", "1", "88.13%", "4.68%", "Strong baseline from 1 example"],
    ["2-shot", "2", "93.11%", "2.64%", "+4.98% gain; variance halved"],
    ["3-shot", "3", "94.86%", "1.85%", "Approaching saturation"],
    ["4-shot", "4", "95.75%", "1.51%", "Near-supervised performance"],
]
col_widths_r = [Inches(1.8), Inches(1.8), Inches(2.2), Inches(1.8), Inches(3.5)]
add_table(slide, Inches(1.1), Inches(1.9), res_data, col_widths_r, font_size=14)

# Highlight boxes
highlights = [
    ("88.13%", "1-shot accuracy", ACCENT_BLUE),
    ("95.75%", "4-shot accuracy", GREEN_ACCENT),
    ("+7.62%", "Improvement 1->4 shot", GOLD),
    ("1.51%", "Lowest std dev (4-shot)", ACCENT_TEAL),
]
box_w = Inches(2.5)
gap = Inches(0.3)
total = len(highlights) * box_w + (len(highlights) - 1) * gap
sx = (SW - total) / 2
for i, (val, label, color) in enumerate(highlights):
    x = sx + i * (box_w + gap)
    add_rounded_rect(slide, x, Inches(4.4), box_w, Inches(1.6), fill_color=CARD_BG)
    add_text_box(slide, x, Inches(4.6), box_w, Inches(0.6),
                 val, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.3), box_w, Inches(0.4),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key insight
add_rounded_rect(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.7), fill_color=CARD_BG)
add_text_box(slide, Inches(1.9), Inches(6.35), Inches(9.5), Inches(0.55),
             "Key Insight: Accuracy improves and variance decreases as more labeled support samples are provided per class. "
             "Self-supervised pretraining enables strong generalization even from a single labeled example.",
             font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 13 – t-SNE VISUALIZATION (with actual result image)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "t-SNE Feature Visualization", number=11)

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.4),
             "2D t-SNE projection of learned test features showing class cluster separation",
             font_size=15, color=LIGHT_GRAY)

# Place the t-SNE result image
if os.path.exists(TSNE_IMG):
    # Center the image
    img_w = Inches(7.5)
    img_h = Inches(5.2)
    img_x = (SW - img_w) / 2
    img_y = Inches(1.8)
    # White border card
    add_rounded_rect(slide, img_x - Inches(0.15), img_y - Inches(0.15),
                     img_w + Inches(0.3), img_h + Inches(0.3), fill_color=WHITE)
    slide.shapes.add_picture(TSNE_IMG, img_x, img_y, img_w, img_h)
else:
    add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                 "[t-SNE visualization image not found]",
                 font_size=24, color=RED_ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.6), Inches(7.1), Inches(12), Inches(0.3),
             "Parameters: perplexity=30  |  Standardized features  |  Fallback: PCA if t-SNE unavailable",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 14 – PROJECT STRUCTURE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Project Structure", number=12)

# File tree
structure_data = [
    ["File / Folder", "Lines", "Role"],
    ["main_TSMMIS.m", "683", "End-to-end orchestrator: config, data, train, eval, visualize"],
    ["utils/DataLoader.m", "284", "Data loading, train/test split, multicrop, augmentation"],
    ["utils/TSMMISConfig.m", "100", "Configuration value-object with all hyperparameters"],
    ["models/ResNet18Encoder.m", "221", "ResNet-18 architecture: conv + 4 residual stages + GAP + FC"],
    ["models/pretrained_encoder.mat", "---", "Saved pretrained model weights"],
    ["training/ContrastiveLoss.m", "99", "L_t&s: cross-view contrastive loss implementation"],
    ["training/MMISLoss.m", "141", "L_MMIS: multi-modal instance similarity regularization"],
    ["training/TSMMISTrainer.m", "262", "OOP training loop with initialize/train/step methods"],
    ["training/trainTSMMIS_GPU.m", "181", "GPU-accelerated training with batch processing"],
    ["evaluation/FineTuner.m", "297", "Few-shot fine-tuning (SVM linear eval + full fine-tune)"],
    ["evaluation/TSMMISEvaluator.m", "285", "Evaluator class with t-SNE visualization"],
    ["evaluation/results/", "---", "Output: evaluation_results.mat, tsne_embedding.mat, tsne_visualization.png"],
    ["data/NEU-CLS/", "---", "6 class folders x 240 images = 1,440 JPG files"],
]
col_widths_s = [Inches(3.2), Inches(0.8), Inches(7.5)]
add_table(slide, Inches(0.9), Inches(1.3), structure_data, col_widths_s, font_size=11)


# =====================================================================
# SLIDE 15 – CONCLUSION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_section_title(slide, "Conclusion", number=13)

# Key achievements
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8), fill_color=CARD_BG)
add_text_box(slide, Inches(0.9), Inches(1.65), Inches(5.2), Inches(0.4),
             "KEY ACHIEVEMENTS", font_size=18, color=GREEN_ACCENT, bold=True)
add_bullet_slide_content(slide, Inches(0.9), Inches(2.2), Inches(5.2), Inches(1.8), [
    "95.75% accuracy with only 4 labeled samples per class",
    "88.13% accuracy with a single labeled sample per class",
    "Self-supervised learning eliminates need for large labeled datasets",
    "Prototype self-training leverages unlabeled data effectively",
], font_size=14, bullet_color=GREEN_ACCENT)

# Technical contributions
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8), fill_color=CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.65), Inches(5.2), Inches(0.4),
             "TECHNICAL CONTRIBUTIONS", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(7.2), Inches(2.2), Inches(5.2), Inches(1.8), [
    "Teacher-Student contrastive learning framework",
    "Multi-Modal Instance Similarity (MMIS) regularization",
    "Combined prototype + instance similarity scoring (75/25)",
    "Confidence-based prototype adaptation from unlabeled data",
], font_size=14, bullet_color=ACCENT_BLUE)

# Future scope
add_rounded_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), fill_color=CARD_BG)
add_text_box(slide, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.4),
             "PRACTICAL IMPACT", font_size=18, color=GOLD, bold=True)
add_bullet_slide_content(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.5), [
    "Enables automated steel surface defect detection in manufacturing environments with minimal labeling effort",
    "Reduces dependency on domain experts for annotation -- only 1-4 labeled examples needed per defect type",
    "Modular MATLAB implementation allows easy extension to new defect types or other industrial inspection tasks",
    "Robust evaluation protocol (100 trials) ensures reliable and reproducible performance estimates",
], font_size=14, bullet_color=GOLD)


# =====================================================================
# SLIDE 16 – THANK YOU
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, 0, 0, SW, Inches(0.06), fill_color=ACCENT_BLUE)
add_shape(slide, 0, SH - Inches(0.06), SW, Inches(0.06), fill_color=ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.8),
             "Thank You", font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_accent_line(slide, Inches(5), Inches(3.2), Inches(3.3), ACCENT_TEAL, Pt(3))

add_text_box(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.6),
             "TSMMIS: Self-Supervised Few-Shot Steel Surface Defect Recognition",
             font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.5),
             "Questions & Discussion",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
